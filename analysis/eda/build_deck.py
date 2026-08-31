"""Build the exploratory-pass slide deck as .pptx and .pdf from one spec.

The same slide specification drives both renderers, so the PDF is a faithful
preview of the PowerPoint geometry rather than a separate design.
"""

from __future__ import annotations

import argparse
from pathlib import Path

W, H = 13.333, 7.5
MARGIN = 0.62

PAPER = "FFFFFF"
DARK = "12222A"
INK = "14202A"
MUTED = "5C6971"
TEAL = "215761"
OCHRE = "A8632A"
RULE = "D6DEDF"
HEAD_FONT = "Cambria"
BODY_FONT = "Calibri"


def slides(figures: Path):
    def fig(name):
        return str(figures/f"{name}.png")

    return [
        {"kind": "title",
         "title": "What the ODyn data can actually support",
         "subtitle": "An exploratory pass over 44 sessions of bulbar dopaminergic imaging, "
                     "run to decide what is worth reporting before the figures are built.",
         "meta": "404,415 unit–odor rows   ·   6 cohorts, 13 mice   ·   "
                 "16-odor panel, pre/post ket-xyl   ·   2026-08-30",
         "notes": "Coverage, signed responses, population geometry and arousal coupling. "
                  "Everything here is descriptive; the conclusions are now implemented in the figure code."},

        {"kind": "stats", "title": "The design decides what can carry weight",
         "takeaway": "The animal-level structure inverts the intuitive ordering of the comparisons.",
         "stats": [("4", "TH mice, each with 10x,\nsuperficial and deep 20x"),
                   ("2", "DAT mice at 10x\nand 2 Thy1 mice"),
                   ("1–11", "trials per odor within\na single session")],
         "bullets": [
             "TH superficial vs deep is a within-mouse contrast in all four TH animals, and pre/post ket-xyl is within-session and within-field. Those two carry the weight.",
             "The headline TH vs DAT line comparison rests on two DAT animals, so it stays descriptive.",
             "Depth is not confounded with line: TH deep 110–200 um, TH superficial 54–103, DAT superficial 53–77.",
         ],
         "notes": "Reported between-line contrasts descriptively with every mouse shown; inferential effort went to the paired designs."},

        {"kind": "figure", "title": "Mineral oil is not a null stimulus",
         "takeaway": "It crosses an excitation threshold in 16–87% of unit–odor pairs, where the nominal rate is 5%.",
         "image": fig("blank"),
         "caption": "Population median response, blank against real odors. The pre-odor baseline is flat, so this is not drift. "
                    "Awake, the blank drives TH glomeruli to roughly two thirds of the real-odor response; at 20x it sits near zero.",
         "notes": "Consequence: thresholds referenced to mineral oil shift between conditions for reasons unrelated to odor coding, "
                  "and the blank response is strongly session-dependent. Detection now references pre-odor excursions instead."},

        {"kind": "figure", "title": "The four-second mean hides most of the structure",
         "takeaway": "Of genuinely bidirectional unit–odor pairs, the signed mean misses 69–88% awake.",
         "image": fig("biphasic"),
         "caption": "One real TH 10x unit. Within the odor window the response goes +1.6, then −4.7, then +2.9 — "
                    "and its signed four-second mean is −0.06 z. Excitation and suppression are now integrated separately.",
         "notes": "Across all pairs the trace method finds a response in 54-96% where the mean-based call finds one in 12-66%."},

        {"kind": "figure", "title": "Anesthesia abolishes suppression, everywhere",
         "takeaway": "The largest and most consistent effect in the dataset, and a within-session contrast.",
         "image": fig("suppression"),
         "caption": "Suppression breadth, mice as sampling units. Every cohort and both compartments move the same way.",
         "notes": "This independently confirms the qualitative read of the group-202 z-score movies."},

        {"kind": "figure", "title": "A depth gradient in excitation/suppression balance",
         "takeaway": "TH deep is the only suppression-dominated population awake — and the most excitatory under ket/xyl.",
         "image": fig("balance"),
         "caption": "Awake somatic balance orders monotonically: DAT superficial +0.67, TH superficial +0.34, TH deep −0.41. "
                    "Awake TH deep somata have a median odor response of −0.19 z.",
         "notes": "The TH superficial vs deep contrast is within-mouse in all four TH animals, so this is the strongest between-population claim available."},

        {"kind": "figure", "title": "How much odor information each population carries",
         "takeaway": "Every population is well above chance — TH deep still classifies 16 odors at four times chance.",
         "image": fig("accuracy"),
         "caption": "Leave-one-trial-out nearest centroid across the whole panel. The 10x/20x gap is expected: "
                    "20x samples a small subset of spatially distributed labeled-line input channels.",
         "notes": "Thy1 is the most informative population in the dataset, which supports using it as the 10x anchor."},

        {"kind": "figure", "title": "What gets mistaken for what",
         "takeaway": "Mixtures are confused with their reciprocal partner, never with their components.",
         "image": fig("confusion"),
         "caption": "Odors ordered so each mixture sits beside its two components. In Thy1 the only off-diagonal mass is "
                    "alpha vs alpha-prime and epsilon vs epsilon-prime; lambda/lambda-prime is fully resolved.",
         "notes": "Ratio is the only failure mode, and it fails for two of the three pairs."},

        {"kind": "figure", "title": "The three mixture families are not alike",
         "takeaway": "Whether ratios are discriminable may be downstream of how a mixture relates to its components.",
         "image": fig("trials"),
         "caption": "Every awake trial against every other, one Thy1 session. Alpha is a novel pattern unlike either component; "
                    "epsilon forms one block with its components; lambda sits between and is the pair whose ratios separate.",
         "notes": "Testable without more data. This is the most hypothesis-generating view of the three."},

        {"kind": "figure", "title": "Ratio separation builds late",
         "takeaway": "Only pair 39/40 separates above chance, and a 0–4 s integral dilutes it.",
         "image": fig("ratio"),
         "caption": "Pair 39/40, awake. Separation peaks in the final second, and for TH deep somata after offset entirely. "
                    "Pairs 17/18 and 31/32 sit at the noise floor in every window. Mixture work is parked until there is more data.",
         "notes": "Permutation p<0.05 rate pooled over windows and cohorts: 39/40 at 0.276, 17/18 at 0.090, 31/32 at 0.055 against 0.05 expected."},

        {"kind": "figure", "title": "Arousal amplifies responses, it does not suppress them",
         "takeaway": "Opposite sign to the NE-suppression hypothesis, and not motion artifact.",
         "image": fig("running"),
         "caption": "Running coupling is flat across the whole pre-odor window, rises 1 s after onset and peaks at 2 s. "
                    "Running occurs in the baseline window too, so an artifact would appear there. It does not.",
         "notes": "Baseline coupling -0.002 against +0.190 in the odor window; only 4 of 57 cases have baseline >= odor, p<0.0001. "
                  "Split by response sign: excited +0.287, suppressed +0.063."},

        {"kind": "figure", "title": "Pupil says the same thing, more weakly",
         "takeaway": "Running dominates, but pupil carries a small independent contribution.",
         "image": fig("pupil"),
         "caption": "Pupil controlling for running +0.055 (p=0.003); running controlling for pupil +0.189 (p<0.0001). "
                    "Carried by dilations, not constrictions. Absent under ket/xyl, where locomotion is zero and pupil quality is best.",
         "notes": "Awake pupil coverage is lost mainly to clipping, up to 41% of samples, and clipping happens when the pupil is large. "
                  "The ROI or fit bounds are worth fixing before pupil carries weight in a figure."},

        {"kind": "bullets", "title": "What changed in the code",
         "takeaway": "The 10x and 20x paths now share one implementation, so the scales are measured identically.",
         "bullets": [
             "Responder calls reference pre-odor excursions, stratified by trial count, replacing the mineral-oil cutoff everywhere.",
             "Excitation and suppression are integrated separately over the waveform, with early/late windows and peak latencies. Biphasic responses are detectable for the first time.",
             "Lifetime sparseness comes from continuous AUC rather than thresholded excess, removing a bias where non-responsive units dropped out in a state-dependent way.",
             "Whole-panel confusion matrices and RDMs are generated alongside the pairwise geometry.",
             "Joins use the source round's database trial ids; states resolve through each file's own levels.",
             "One crossnobis implementation, with a permutation null that flags underpowered rather than emitting a z score.",
         ],
         "notes": "155 tests pass across the repo."},

        {"kind": "closing", "title": "Open questions",
         "bullets": [
             "What changed around 2026-08-03? It is the largest single source of 10x session variance, and it survives within line.",
             "Lagged and FIR models of the arousal coupling; respiration-phase regressors.",
             "The z-stack soma-size by depth model.",
             "Awake pupil clipping, before pupil carries weight in a figure.",
             "Continuous inter-trial auxiliary data, if the spontaneous-event analysis is worth recovering.",
         ],
         "notes": "Mixture work is parked pending more data."},
    ]


def _fit(image, box):
    """Scale an image into a box, preserving aspect, centred."""
    from PIL import Image

    x, y, width, height = box
    with Image.open(image) as handle:
        aspect = handle.width/handle.height
    if width/height > aspect:
        drawn_h, drawn_w = height, height*aspect
    else:
        drawn_w, drawn_h = width, width/aspect
    return x + (width-drawn_w)/2, y + (height-drawn_h)/2, drawn_w, drawn_h


def build_pptx(spec, path):
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    def rgb(value):
        return RGBColor.from_string(value)

    presentation = Presentation()
    presentation.slide_width = Inches(W)
    presentation.slide_height = Inches(H)
    blank = presentation.slide_layouts[6]

    def textbox(slide, text, box, *, size, colour, font=BODY_FONT, bold=False,
                italic=False, spacing=1.0, bullet=False):
        x, y, width, height = box
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width),
                                         Inches(height))
        frame = shape.text_frame
        frame.word_wrap = True
        frame.margin_left = frame.margin_right = 0
        frame.margin_top = frame.margin_bottom = 0
        items = text if isinstance(text, list) else [text]
        for index, item in enumerate(items):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.line_spacing = spacing
            if index:
                paragraph.space_before = Pt(9)
            run = paragraph.add_run()
            run.text = ("•  " + item) if bullet else item
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.name = font
            run.font.color.rgb = rgb(colour)
        return shape

    def background(slide, colour):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = rgb(colour)

    for item in spec:
        slide = presentation.slides.add_slide(blank)
        kind = item["kind"]
        dark = kind in ("title", "closing")
        background(slide, DARK if dark else PAPER)
        head = PAPER if dark else INK
        soft = "9FB0B6" if dark else MUTED

        if kind == "title":
            textbox(slide, item["title"], (MARGIN, 2.05, W-2*MARGIN-1.4, 1.9),
                    size=40, colour=head, font=HEAD_FONT, bold=True, spacing=1.05)
            textbox(slide, item["subtitle"], (MARGIN, 4.15, 8.6, 1.2),
                    size=15, colour=soft, spacing=1.25)
            textbox(slide, item["meta"], (MARGIN, 6.35, W-2*MARGIN, .5),
                    size=11, colour=soft)
        elif kind == "closing":
            textbox(slide, item["title"], (MARGIN, .85, W-2*MARGIN, .9),
                    size=34, colour=head, font=HEAD_FONT, bold=True)
            textbox(slide, item["bullets"], (MARGIN, 2.1, W-2*MARGIN-1.0, 4.6),
                    size=15, colour=soft, spacing=1.3, bullet=True)
        else:
            textbox(slide, item["title"], (MARGIN, .48, W-2*MARGIN, .78),
                    size=28, colour=head, font=HEAD_FONT, bold=True)
            if item.get("takeaway"):
                textbox(slide, item["takeaway"], (MARGIN, 1.24, W-2*MARGIN-.4, .5),
                        size=14, colour=TEAL, italic=True)
            if kind == "figure":
                caption_h = 0.86
                box = (MARGIN, 1.92, W-2*MARGIN, H-1.92-caption_h-.45)
                x, y, width, height = _fit(item["image"], box)
                slide.shapes.add_picture(item["image"], Inches(x), Inches(y),
                                         Inches(width), Inches(height))
                textbox(slide, item["caption"], (MARGIN, H-caption_h-.34,
                                                 W-2*MARGIN, caption_h),
                        size=10.5, colour=soft, spacing=1.22)
            elif kind == "stats":
                y = 2.0
                width = (W-2*MARGIN-1.0)/3
                for index, (value, label) in enumerate(item["stats"]):
                    x = MARGIN + index*(width+.5)
                    textbox(slide, value, (x, y, width, 1.0), size=52,
                            colour=OCHRE, font=HEAD_FONT, bold=True)
                    textbox(slide, label.replace("\n", " "), (x, y+1.02, width, .8),
                            size=11.5, colour=soft, spacing=1.2)
                textbox(slide, item["bullets"], (MARGIN, 4.35, W-2*MARGIN-.6, 2.5),
                        size=13.5, colour=INK, spacing=1.3, bullet=True)
            elif kind == "bullets":
                textbox(slide, item["bullets"], (MARGIN, 2.0, W-2*MARGIN-.6, 4.8),
                        size=14, colour=INK, spacing=1.3, bullet=True)
        if item.get("notes"):
            slide.notes_slide.notes_text_frame.text = item["notes"]

    presentation.save(path)
    return path


def build_pdf(spec, path, png_dir=None):
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from matplotlib.backends.backend_pdf import PdfPages
    import matplotlib.image as mpimg

    def hexc(value):
        return "#" + value

    with PdfPages(path) as pdf:
        for item in spec:
            kind = item["kind"]
            dark = kind in ("title", "closing")
            ground = hexc(DARK if dark else PAPER)
            head = hexc(PAPER if dark else INK)
            soft = hexc("9FB0B6" if dark else MUTED)
            figure = plt.figure(figsize=(W, H), facecolor=ground)
            ax = figure.add_axes([0, 0, 1, 1]); ax.set_axis_off()
            ax.add_patch(plt.Rectangle((0, 0), 1, 1, transform=ax.transAxes,
                                       facecolor=ground, zorder=0))

            def text(value, x, y, size, colour, *, weight="normal",
                     style="normal", family="sans-serif", width=None, va="top"):
                ax.text(x/W, 1-y/H, value, transform=ax.transAxes, fontsize=size,
                        color=colour, weight=weight, style=style, family=family,
                        va=va, ha="left", wrap=True,
                        bbox=dict(boxstyle="square,pad=0", fc="none", ec="none"))

            def wrapped(lines, x, y, size, colour, *, step, bullet=False,
                        chars=118):
                import textwrap
                cursor = y
                for line in lines:
                    for part in textwrap.wrap(("•  " if bullet else "") + line,
                                              chars):
                        text(part, x, cursor, size, colour)
                        cursor += step
                        chars_first = chars
                    cursor += step*.45
                return cursor

            if kind == "title":
                import textwrap
                cursor = 2.35
                for part in textwrap.wrap(item["title"], 38):
                    text(part, MARGIN, cursor, 34, head, weight="bold",
                         family="serif")
                    cursor += .72
                wrapped([item["subtitle"]], MARGIN, 4.35, 13, soft, step=.32, chars=92)
                text(item["meta"], MARGIN, 6.5, 9.5, soft)
            elif kind == "closing":
                text(item["title"], MARGIN, 1.0, 29, head, weight="bold",
                     family="serif")
                wrapped(item["bullets"], MARGIN, 2.25, 12.5, soft, step=.30,
                        bullet=True, chars=104)
            else:
                text(item["title"], MARGIN, .62, 24, head, weight="bold",
                     family="serif")
                if item.get("takeaway"):
                    wrapped([item["takeaway"]], MARGIN, 1.36, 12.5,
                            hexc(TEAL), step=.28, chars=120)
                if kind == "figure":
                    box = (MARGIN, 1.95, W-2*MARGIN, H-1.95-1.3)
                    x, y, width, height = _fit(item["image"], box)
                    figure.figimage(mpimg.imread(item["image"]), 0, 0, zorder=1,
                                    resize=False, alpha=0)
                    inset = figure.add_axes([x/W, 1-(y+height)/H, width/W, height/H])
                    inset.imshow(mpimg.imread(item["image"]))
                    inset.set_axis_off()
                    wrapped([item["caption"]], MARGIN, H-1.12, 9.5, soft,
                            step=.24, chars=155)
                elif kind == "stats":
                    width = (W-2*MARGIN-1.0)/3
                    for index, (value, label) in enumerate(item["stats"]):
                        x = MARGIN + index*(width+.5)
                        text(value, x, 2.15, 44, hexc(OCHRE), weight="bold",
                             family="serif")
                        for offset, part in enumerate(label.split("\n")):
                            text(part, x, 3.25+offset*.28, 10.5, soft)
                    wrapped(item["bullets"], MARGIN, 4.5, 12, head, step=.28,
                            bullet=True, chars=122)
                elif kind == "bullets":
                    wrapped(item["bullets"], MARGIN, 2.15, 12.5, head, step=.29,
                            bullet=True, chars=118)
            pdf.savefig(figure, facecolor=ground)
            if png_dir is not None:
                png_dir.mkdir(parents=True, exist_ok=True)
                figure.savefig(png_dir/f"slide-{spec.index(item)+1:02d}.png",
                               dpi=110, facecolor=ground)
            plt.close(figure)
    return path


def main(argv=None):
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--figures", type=Path, required=True)
    parser.add_argument("--output-dir", type=Path, required=True)
    parser.add_argument("--stem", default="odyn_exploratory_pass")
    parser.add_argument("--png-dir", type=Path, default=None,
                        help="also write one PNG per slide, for review")
    args = parser.parse_args(argv)
    args.output_dir.mkdir(parents=True, exist_ok=True)
    spec = slides(args.figures)
    pptx = build_pptx(spec, args.output_dir/f"{args.stem}.pptx")
    pdf = build_pdf(spec, args.output_dir/f"{args.stem}_slides.pdf",
                    png_dir=args.png_dir)
    print(f"{len(spec)} slides")
    for path in (pptx, pdf):
        print(f"  {path}  ({path.stat().st_size/1e6:.2f} MB)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
